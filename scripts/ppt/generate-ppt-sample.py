#!/usr/bin/env python3
"""Generate the first-five-event BenchCouncil AI100 PowerPoint sample."""

from __future__ import annotations

import html
import json
import re
import shutil
import subprocess
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / ".tmp" / "ppt-sample"
ASSET_CACHE = OUT / "assets-v2"
STORYLINE_PATH = ROOT / "archive" / "storylines" / "bench-council-ai100.json"
FIGURES_PATH = ROOT / "archive" / "figures" / "figures.json"

PAPER = RGBColor(246, 243, 235)
INK = RGBColor(26, 35, 39)
INK_SOFT = RGBColor(73, 84, 87)
TEAL = RGBColor(13, 111, 105)
TEAL_PALE = RGBColor(221, 235, 229)
VERMILION = RGBColor(190, 66, 50)
OCHRE = RGBColor(189, 139, 54)
HAIRLINE = RGBColor(198, 195, 184)
WHITE = RGBColor(255, 255, 255)

FONT_ZH = "Microsoft YaHei"
FONT_LATIN = "Aptos"

EVENT_IDS = [
    "1950-turing-test",
    "1971-complexity-theory",
    "1971-vc-theory",
    "1956-logic-theorist",
    "1958-wangs-algorithm",
]


def strip_html(value: str | None) -> str:
    value = value or ""
    value = re.sub(r"<br\s*/?>", "\n", value, flags=re.I)
    value = re.sub(r"</p\s*>", "\n\n", value, flags=re.I)
    value = re.sub(r"<[^>]+>", "", value)
    return html.unescape(re.sub(r"\n{3,}", "\n\n", value)).strip()


def load_json(path: Path):
    return json.loads(path.read_text(encoding="utf-8"))


def load_figures() -> dict[str, dict]:
    return {item["id"]: item for item in load_json(FIGURES_PATH)}


def event_bundle(event_id: str) -> dict:
    base = ROOT / "archive" / "events" / event_id
    event = load_json(base / "event.json")
    event["_assets"] = load_json(base / "assets.json")
    return event


def solid_rect(slide, x, y, w, h, color, line=None, radius=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = line or color
    shape.line.width = Pt(0.8)
    return shape


def add_rule(slide, x, y, w, color=HAIRLINE, thickness=1.2):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.012))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.color.rgb = color
    line.line.width = Pt(thickness)
    return line


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=22,
    color=INK,
    bold=False,
    font=FONT_ZH,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    fit=True,
    margin=0.02,
    line_spacing=1.15,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    if fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_paragraph_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=17,
    color=INK,
    bold=False,
    line_spacing=1.22,
    paragraph_spacing=7,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for index, paragraph_text in enumerate(part for part in text.split("\n\n") if part.strip()):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.line_spacing = line_spacing
        paragraph.space_after = Pt(paragraph_spacing)
        run = paragraph.add_run()
        run.text = paragraph_text.strip()
        run.font.name = FONT_ZH
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_label(slide, text, x, y, w, color=TEAL):
    add_text(slide, text, x, y, w, 0.24, 10, color, True, FONT_ZH, fit=False)


def add_footer(slide, page, total, event_id=None):
    add_rule(slide, 0.65, 7.15, 12.03, HAIRLINE, 0.8)
    left = "AI 顶尖成就图谱（BenchCouncil） · 前 5 个事件样例"
    if event_id:
        left += f"  ·  {event_id}"
    add_text(slide, left, 0.67, 7.23, 10.9, 0.15, 7.5, INK_SOFT, fit=False)
    add_text(slide, f"{page:02d} / {total:02d}", 11.72, 7.22, 0.95, 0.15, 7.5, INK_SOFT, align=PP_ALIGN.RIGHT, fit=False)


def add_header(slide, kicker, title, meta, page, total, event_id=None):
    add_text(slide, kicker, 0.68, 0.25, 4.6, 0.20, 8.5, TEAL, True, fit=False)
    add_text(slide, title, 0.65, 0.51, 8.55, 0.46, 23, INK, True)
    add_text(slide, meta, 9.25, 0.57, 3.42, 0.24, 9, INK_SOFT, align=PP_ALIGN.RIGHT)
    add_rule(slide, 0.65, 1.10, 12.03, HAIRLINE)
    add_footer(slide, page, total, event_id)


def find_asset(event, roles):
    presentation = event.get("defaultPresentation", {})
    selected_ids = []
    overview_id = presentation.get("overviewImageAssetId")
    if overview_id:
        selected_ids.append(overview_id)
    selected_ids.extend(presentation.get("assetIds", []))
    assets_by_id = {asset.get("id"): asset for asset in event.get("_assets", [])}
    for role in roles:
        for asset_id in selected_ids:
            asset = assets_by_id.get(asset_id)
            if not asset:
                continue
            if asset.get("role") != role:
                continue
            path = asset.get("path")
            if path and (ROOT / path).exists():
                return ROOT / path
    return None


def rasterize(path: Path | None) -> Path | None:
    if not path:
        return None
    if path.suffix.lower() not in {".svg", ".svgz"}:
        return path
    if not shutil.which("convert"):
        raise RuntimeError(
            "SVG 图片需要 ImageMagick 的 convert 命令；"
            "请先安装 ImageMagick，详见 scripts/ppt/README.md。"
        )
    out = ASSET_CACHE / f"{path.stem}.png"
    if not out.exists() or out.stat().st_mtime < path.stat().st_mtime:
        subprocess.run(
            ["convert", "-background", "none", "-density", "180", str(path), str(out)],
            check=True,
        )
    return out


def add_picture_crop(slide, path: Path | None, x, y, w, h, contain=False):
    path = rasterize(path)
    if not path:
        return None
    with Image.open(path) as image:
        image_w, image_h = image.size
    if contain:
        scale = min(w / image_w, h / image_h)
        draw_w, draw_h = image_w * scale, image_h * scale
        draw_x = x + (w - draw_w) / 2
        draw_y = y + (h - draw_h) / 2
        return slide.shapes.add_picture(str(path), Inches(draw_x), Inches(draw_y), Inches(draw_w), Inches(draw_h))
    image_ratio = image_w / image_h
    box_ratio = w / h
    picture = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    if image_ratio > box_ratio:
        visible = box_ratio / image_ratio
        crop = (1 - visible) / 2
        picture.crop_left = crop
        picture.crop_right = crop
    else:
        visible = image_ratio / box_ratio
        crop = (1 - visible) / 2
        picture.crop_top = crop
        picture.crop_bottom = crop
    return picture


def figure_lines(event, figure_map):
    lines = []
    for figure in event.get("figures", [])[:4]:
        entry = figure_map.get(figure.get("figureId"), {})
        name = entry.get("name", {}).get("zh") or figure.get("figureId", "")
        role = figure.get("role", {}).get("zh", "")
        lines.append((name, role))
    return lines


def get_section(event, candidates):
    sections = event.get("defaultPresentation", {}).get("commentarySections", [])
    for candidate in candidates:
        for section in sections:
            if section.get("id") == candidate:
                return strip_html(section.get("html", {}).get("zh"))
    return ""


def add_overview_slide(prs, event, position, page, total, figure_map):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    year = str(event.get("year", ""))
    title = event["title"]["zh"]
    add_header(slide, f"事件 {position:02d} · 历史概览", title, f"{year} · BenchCouncil AI100", page, total, event["id"])

    portrait = find_asset(event, ["portrait", "source-card"])
    solid_rect(slide, 0.66, 1.36, 3.72, 5.55, WHITE, HAIRLINE)
    add_picture_crop(slide, portrait, 0.84, 1.54, 3.36, 4.24)
    add_text(slide, year, 0.92, 6.02, 2.35, 0.50, 24, VERMILION, True, FONT_LATIN)
    add_text(slide, "AI HISTORY", 3.02, 6.17, 1.00, 0.18, 8, INK_SOFT, True, FONT_LATIN, PP_ALIGN.RIGHT, fit=False)

    add_label(slide, "事件详述", 4.80, 1.43, 2.0)
    summary = strip_html(
        event.get("defaultPresentation", {}).get("displayDescription", {}).get("zh")
        or event.get("description", {}).get("zh")
    )
    add_paragraph_text(slide, summary, 4.78, 1.79, 7.72, 3.20, 16.5, INK, line_spacing=1.23, paragraph_spacing=9)
    add_rule(slide, 4.80, 5.15, 7.67, TEAL, 1.4)

    add_label(slide, "关键人物", 4.80, 5.39, 2.0)
    figures = figure_lines(event, figure_map)
    row_y = 5.71
    for idx, (name, role) in enumerate(figures):
        col = idx % 2
        row = idx // 2
        x = 4.80 + col * 3.87
        y = row_y + row * 0.62
        solid_rect(slide, x, y + 0.04, 0.07, 0.43, TEAL if idx == 0 else OCHRE)
        add_text(slide, name, x + 0.17, y, 3.43, 0.25, 13, INK, True)
        add_text(slide, role, x + 0.17, y + 0.27, 3.43, 0.22, 9.5, INK_SOFT)

    topics = " / ".join(event.get("topics", [])) or "AI 历史"
    add_text(slide, topics, 4.80, 6.78, 7.55, 0.16, 8.5, TEAL, True, FONT_LATIN, fit=False)
    return slide


def add_mechanism_slide(prs, event, position, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    title = event["title"]["zh"]
    add_header(slide, f"事件 {position:02d} · 技术解释", title, "核心机制与长期影响", page, total, event["id"])

    explainer = find_asset(event, ["algorithm-explainer", "architecture-explainer"])
    solid_rect(slide, 0.66, 1.36, 6.35, 5.55, WHITE, HAIRLINE)
    add_picture_crop(slide, explainer, 0.93, 1.62, 5.81, 4.72, contain=True)
    add_text(slide, "ARCHIVE EXPLAINER", 0.94, 6.54, 2.25, 0.16, 7.5, INK_SOFT, True, FONT_LATIN, fit=False)

    solid_rect(slide, 7.36, 1.36, 5.31, 2.60, TEAL_PALE, TEAL_PALE)
    add_label(slide, "核心思想", 7.70, 1.68, 1.5, TEAL)
    core = get_section(event, ["reasoning-logic", "reduction-logic", "capacity-logic", "evaluation-logic", "pattern-logic", "core-idea"])
    add_text(slide, core, 7.68, 2.06, 4.55, 1.48, 16.5, INK, line_spacing=1.25)

    solid_rect(slide, 7.36, 4.24, 5.31, 2.67, WHITE, HAIRLINE)
    add_label(slide, "长期影响", 7.70, 4.57, 1.5, VERMILION)
    legacy = get_section(event, ["long-term-legacy"])
    add_text(slide, legacy, 7.68, 4.96, 4.55, 1.52, 15.5, INK_SOFT, line_spacing=1.25)
    return slide


def add_cover(prs, storyline, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    solid_rect(slide, 0, 0, 0.22, 7.5, VERMILION)
    solid_rect(slide, 8.82, 0, 4.51, 7.5, TEAL)
    add_text(slide, "AI HISTORY · BENCHCOUNCIL", 0.82, 0.86, 6.3, 0.28, 12, TEAL, True, FONT_LATIN, fit=False)
    cover_title = storyline["title"]["zh"].replace("（BenchCouncil）", "")
    add_text(slide, cover_title, 0.78, 1.52, 7.35, 0.72, 36, INK, True, line_spacing=1.05)
    add_text(slide, "BenchCouncil", 0.80, 2.31, 5.20, 0.44, 24, INK, True, FONT_LATIN, fit=False)
    add_text(slide, "前 5 个事件 · PPT 样例", 0.82, 3.16, 5.8, 0.48, 22, VERMILION, True)
    add_rule(slide, 0.82, 3.92, 6.86, HAIRLINE)
    add_text(slide, "每个事件两页：历史概览 / 技术解释", 0.82, 4.25, 6.45, 0.42, 17, INK_SOFT)
    add_text(slide, "Archive JSON 驱动 · 可编辑 PPTX", 0.82, 4.90, 5.8, 0.30, 11, INK_SOFT, True)
    add_text(slide, "05", 9.32, 1.15, 2.65, 1.12, 58, WHITE, True, FONT_LATIN)
    add_text(slide, "EVENTS", 9.38, 2.38, 2.10, 0.30, 14, WHITE, True, FONT_LATIN, fit=False)
    add_text(slide, "图灵测试\nNP 完全性\nVC 理论\n逻辑理论家\n王氏算法", 9.38, 3.25, 2.72, 2.15, 19, WHITE, True, line_spacing=1.38)
    add_footer(slide, 1, total)


def add_contents(prs, storyline, events, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_header(slide, "样例目录", "五个事件，两种页面", "历史概览 + 技术解释", 2, total)
    add_text(slide, "一个事件不限于一页，但本轮只保留最有用的两页。", 0.72, 1.90, 8.8, 0.46, 21, INK_SOFT)
    y = 2.72
    for idx, event in enumerate(events, 1):
        color = TEAL if idx % 2 else VERMILION
        add_text(slide, f"{idx:02d}", 0.75, y, 0.72, 0.42, 20, color, True, FONT_LATIN, fit=False)
        add_text(slide, event["title"]["zh"], 1.56, y - 0.01, 4.55, 0.42, 18, INK, True)
        add_text(slide, str(event.get("year", "")), 6.36, y + 0.03, 1.50, 0.30, 13, INK_SOFT, True, FONT_LATIN)
        add_text(slide, f"第 {3 + (idx - 1) * 2}–{4 + (idx - 1) * 2} 页", 9.66, y + 0.03, 2.35, 0.30, 12, INK_SOFT, align=PP_ALIGN.RIGHT)
        add_rule(slide, 1.55, y + 0.58, 10.45, HAIRLINE, 0.6)
        y += 0.82
    add_footer(slide, 2, total)


def validate_layout(prs):
    issues = []
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    for slide_no, slide in enumerate(prs.slides, 1):
        for shape_no, shape in enumerate(slide.shapes, 1):
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > slide_w or shape.top + shape.height > slide_h:
                issues.append(f"slide {slide_no} shape {shape_no} exceeds slide bounds")
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                if shape.width < Inches(0.18) or shape.height < Inches(0.12):
                    issues.append(f"slide {slide_no} text shape {shape_no} is too small")
    if issues:
        raise RuntimeError("\n".join(issues))


def main():
    OUT.mkdir(parents=True, exist_ok=True)
    ASSET_CACHE.mkdir(parents=True, exist_ok=True)
    storyline = load_json(STORYLINE_PATH)
    figure_map = load_figures()
    events = [event_bundle(event_id) for event_id in EVENT_IDS]

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    total = 2 + len(events) * 2
    index = {"storylineId": storyline["id"], "title": storyline["title"]["zh"], "slides": []}

    add_cover(prs, storyline, total)
    add_contents(prs, storyline, events, total)
    for position, event in enumerate(events, 1):
        start = len(prs.slides) + 1
        add_overview_slide(prs, event, position, start, total, figure_map)
        add_mechanism_slide(prs, event, position, start + 1, total)
        index["slides"].append(
            {
                "eventId": event["id"],
                "title": event["title"]["zh"],
                "slideStart": start,
                "slideEnd": start + 1,
            }
        )

    validate_layout(prs)
    out_pptx = OUT / "ai-history-bench-council-first-5-sample-v3.pptx"
    out_index = OUT / "slide-event-index-v3.json"
    prs.save(out_pptx)
    out_index.write_text(json.dumps(index, ensure_ascii=False, indent=2), encoding="utf-8")
    print(out_pptx)
    print(out_index)


if __name__ == "__main__":
    main()
