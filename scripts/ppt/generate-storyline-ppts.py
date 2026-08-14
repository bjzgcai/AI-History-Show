#!/usr/bin/env python3
"""Generate one editable PowerPoint deck per AI History Show storyline."""

from __future__ import annotations

import argparse
import hashlib
import html
import json
import math
import re
import shutil
import subprocess
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
STORYLINE_DIR = ROOT / "archive" / "storylines"
EVENT_DIR = ROOT / "archive" / "events"
FIGURES_PATH = ROOT / "archive" / "figures" / "figures.json"
DEFAULT_OUTPUT = ROOT / "exports" / "ai-history-ppt"

PAPER = RGBColor(246, 243, 235)
INK = RGBColor(26, 35, 39)
INK_SOFT = RGBColor(73, 84, 87)
TEAL = RGBColor(13, 111, 105)
TEAL_PALE = RGBColor(221, 235, 229)
VERMILION = RGBColor(190, 66, 50)
OCHRE = RGBColor(189, 139, 54)
HAIRLINE = RGBColor(198, 195, 184)
WHITE = RGBColor(255, 255, 255)

FONT_ZH = "Microsoft YaHei"
FONT_LATIN = "Aptos"
TOC_EVENTS_PER_PAGE = 10

STORYLINES = [
    {
        "id": "bench-council-ai100",
        "filename": "01-ai-achievement-map-benchcouncil.pptx",
        "indexFilename": "01-ai-achievement-map-benchcouncil-index.json",
        "coverTag": "AI HISTORY · BENCHCOUNCIL",
    },
    {
        "id": "deep-learning",
        "filename": "02-connectionism-seventy-years.pptx",
        "indexFilename": "02-connectionism-seventy-years-index.json",
        "coverTag": "AI HISTORY · CONNECTIONISM",
    },
    {
        "id": "gaming-ai",
        "filename": "03-gaming-ai-evolution.pptx",
        "indexFilename": "03-gaming-ai-evolution-index.json",
        "coverTag": "AI HISTORY · GAME INTELLIGENCE",
    },
    {
        "id": "humanistic-cycle",
        "filename": "04-humanistic-emotional-cycles.pptx",
        "indexFilename": "04-humanistic-emotional-cycles-index.json",
        "coverTag": "AI HISTORY · HUMANISTIC CYCLES",
    },
]

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".svg", ".svgz"}
OVERVIEW_ROLES = [
    "portrait",
    "source-card",
    "supporting-portrait",
    "team-portrait",
    "team-photo",
    "historical-photo",
    "hero-image",
    "venue-photo",
    "project-identity",
]
EXPLAINER_ROLES = [
    "algorithm-explainer",
    "architecture-explainer",
    "annual-achievement-explainer",
    "architecture-diagram",
    "historical-diagram",
    "paper-figure",
    "research-result",
    "game-evolution-poster",
    "paper-case-poster",
    "game-record-image",
    "game-comparison-image",
    "game-analysis-image",
    "gameplay-image",
    "historical-reconstruction",
    "paper-page",
    "primary-source",
    "paper-record",
    "supporting-image",
    "hero-image",
]
FILL_ROLES = {
    "portrait",
    "source-card",
    "supporting-portrait",
    "team-portrait",
    "team-photo",
    "historical-photo",
    "venue-photo",
}


def strip_html(value: str | None) -> str:
    value = value or ""
    value = re.sub(r"<br\s*/?>", "\n", value, flags=re.I)
    value = re.sub(r"</p\s*>", "\n\n", value, flags=re.I)
    value = re.sub(r"<li\s*>", "", value, flags=re.I)
    value = re.sub(r"</li\s*>", "\n", value, flags=re.I)
    value = re.sub(r"<[^>]+>", "", value)
    return html.unescape(re.sub(r"\n{3,}", "\n\n", value)).strip()


def load_json(path: Path):
    return json.loads(path.read_text(encoding="utf-8"))


def localized(value, language="zh", default=""):
    if isinstance(value, dict):
        return value.get(language) or value.get("en") or default
    return value or default


def deck_title(storyline) -> str:
    return localized(storyline.get("title")).replace("（BenchCouncil）", " · BenchCouncil")


def load_figures() -> dict[str, dict]:
    return {item["id"]: item for item in load_json(FIGURES_PATH)}


def load_storyline(storyline_id: str) -> dict:
    storyline = load_json(STORYLINE_DIR / f"{storyline_id}.json")
    storyline["events"] = sorted(
        [event for event in storyline.get("events", []) if event.get("enabled", True)],
        key=lambda event: (event.get("order", 0), event.get("eventId", "")),
    )
    return storyline


def load_event(event_id: str) -> dict:
    base = EVENT_DIR / event_id
    event = load_json(base / "event.json")
    event["_assets"] = load_json(base / "assets.json")
    return event


def solid_rect(slide, x, y, w, h, color, line=None, radius=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = line or color
    shape.line.width = Pt(0.8)
    return shape


def add_rule(slide, x, y, w, color=HAIRLINE, thickness=1.2):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.012))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.color.rgb = color
    line.line.width = Pt(thickness)
    return line


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=22,
    color=INK,
    bold=False,
    font=FONT_ZH,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    fit=True,
    margin=0.02,
    line_spacing=1.15,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    if fit:
        frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = str(text or "")
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_paragraph_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=17,
    color=INK,
    line_spacing=1.22,
    paragraph_spacing=7,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    parts = [part.strip() for part in text.split("\n\n") if part.strip()]
    if not parts:
        parts = ["Archive 暂无可展示的详细说明。"]
    for index, paragraph_text in enumerate(parts):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.line_spacing = line_spacing
        paragraph.space_after = Pt(paragraph_spacing)
        run = paragraph.add_run()
        run.text = paragraph_text
        run.font.name = FONT_ZH
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_label(slide, text, x, y, w, color=TEAL):
    add_text(slide, text, x, y, w, 0.24, 10, color, True, FONT_ZH, fit=False)


def add_footer(slide, storyline_title, page, total, event_id=None):
    add_rule(slide, 0.65, 7.15, 12.03, HAIRLINE, 0.8)
    left = storyline_title
    if event_id:
        left += f"  ·  {event_id}"
    add_text(slide, left, 0.67, 7.23, 10.9, 0.15, 7.5, INK_SOFT, fit=False)
    add_text(
        slide,
        f"{page:03d} / {total:03d}",
        11.55,
        7.22,
        1.12,
        0.15,
        7.5,
        INK_SOFT,
        align=PP_ALIGN.RIGHT,
        fit=False,
    )


def add_header(slide, storyline_title, kicker, title, meta, page, total, event_id=None):
    add_text(slide, kicker, 0.68, 0.25, 4.6, 0.20, 8.5, TEAL, True, fit=False)
    add_text(slide, title, 0.65, 0.51, 8.55, 0.46, 23, INK, True)
    add_text(slide, meta, 9.25, 0.57, 3.42, 0.24, 9, INK_SOFT, align=PP_ALIGN.RIGHT)
    add_rule(slide, 0.65, 1.10, 12.03, HAIRLINE)
    add_footer(slide, storyline_title, page, total, event_id)


def existing_image_assets(event) -> list[dict]:
    assets = []
    for asset in event.get("_assets", []):
        path = asset.get("path")
        if not path or Path(path).suffix.lower() not in IMAGE_EXTENSIONS:
            continue
        absolute = ROOT / path
        if absolute.exists():
            item = dict(asset)
            item["_path"] = absolute
            assets.append(item)
    return assets


def ordered_image_assets(event) -> list[dict]:
    assets = existing_image_assets(event)
    by_id = {asset.get("id"): asset for asset in assets}
    presentation = event.get("defaultPresentation", {})
    order = []
    overview_id = presentation.get("overviewImageAssetId")
    if overview_id:
        order.append(overview_id)
    order.extend(presentation.get("assetIds", []))
    ordered = []
    used = set()
    for asset_id in order:
        asset = by_id.get(asset_id)
        if asset and asset_id not in used:
            ordered.append(asset)
            used.add(asset_id)
    return ordered


def choose_asset(event, preferred_roles, excluded_ids=None) -> tuple[dict | None, bool]:
    excluded_ids = set(excluded_ids or [])
    ordered = ordered_image_assets(event)
    for role in preferred_roles:
        for asset in ordered:
            if asset.get("id") not in excluded_ids and asset.get("role") == role:
                return asset, False
    for asset in ordered:
        if asset.get("id") not in excluded_ids:
            return asset, True
    return None, True


def rasterize(path: Path | None, cache_dir: Path) -> Path | None:
    if not path:
        return None
    suffix = path.suffix.lower()
    digest = hashlib.sha1(str(path).encode("utf-8")).hexdigest()[:12]
    output = cache_dir / f"{path.stem}-{digest}.png"
    if output.exists() and output.stat().st_mtime >= path.stat().st_mtime:
        return output
    if suffix in {".svg", ".svgz"}:
        if not shutil.which("convert"):
            raise RuntimeError(
                "SVG 图片需要 ImageMagick 的 convert 命令；"
                "请先安装 ImageMagick，详见 scripts/ppt/README.md。"
            )
        subprocess.run(
            ["convert", "-background", "none", "-density", "180", str(path), str(output)],
            check=True,
        )
    else:
        with Image.open(path) as source_image:
            source_format = (source_image.format or "").upper()
        if source_format in {"BMP", "GIF", "JPEG", "PNG", "TIFF", "WMF"}:
            return path
        with Image.open(path) as image:
            image.seek(0)
            image.convert("RGBA").save(output)
    return output


def add_picture(slide, asset, cache_dir, x, y, w, h, contain=False):
    if not asset:
        add_text(slide, "暂无可展示图片", x, y + h / 2 - 0.2, w, 0.4, 15, INK_SOFT, align=PP_ALIGN.CENTER)
        return None
    path = rasterize(asset.get("_path"), cache_dir)
    with Image.open(path) as image:
        image_w, image_h = image.size
    if contain:
        scale = min(w / image_w, h / image_h)
        draw_w, draw_h = image_w * scale, image_h * scale
        draw_x = x + (w - draw_w) / 2
        draw_y = y + (h - draw_h) / 2
        return slide.shapes.add_picture(str(path), Inches(draw_x), Inches(draw_y), Inches(draw_w), Inches(draw_h))
    image_ratio = image_w / image_h
    box_ratio = w / h
    picture = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    if image_ratio > box_ratio:
        visible = box_ratio / image_ratio
        crop = (1 - visible) / 2
        picture.crop_left = crop
        picture.crop_right = crop
    else:
        visible = image_ratio / box_ratio
        crop = (1 - visible) / 2
        picture.crop_top = crop
        picture.crop_bottom = crop
    return picture


def figure_lines(event, figure_map):
    lines = []
    for figure in event.get("figures", [])[:4]:
        entry = figure_map.get(figure.get("figureId"), {})
        name = localized(entry.get("name"), default=figure.get("figureId", ""))
        role = localized(figure.get("role"))
        lines.append((name, role))
    return lines


def sections(event) -> list[dict]:
    return event.get("defaultPresentation", {}).get("commentarySections", [])


def section_text(section) -> str:
    return strip_html(localized(section.get("html")))


def exact_section(event, candidates) -> str:
    for candidate in candidates:
        for section in sections(event):
            if section.get("id") == candidate and section_text(section):
                return section_text(section)
    return ""


def detailed_description(event) -> str:
    presentation = event.get("defaultPresentation", {})
    return strip_html(
        localized(presentation.get("displayDescription"))
        or localized(event.get("description"))
        or localized(event.get("summary"))
    )


def content_sections(event) -> tuple[str, str]:
    available = [(section, section_text(section)) for section in sections(event) if section_text(section)]
    core = exact_section(
        event,
        [
            "reasoning-logic",
            "reduction-logic",
            "capacity-logic",
            "evaluation-logic",
            "pattern-logic",
            "core-idea",
        ],
    )
    legacy = exact_section(event, ["long-term-legacy"])
    if not core:
        for section, text in available:
            label = localized(section.get("label"))
            if any(token in label for token in ("核心", "原理", "机制", "思想")):
                core = text
                break
    if not core and available:
        core = available[0][1]
    if not legacy:
        for section, text in reversed(available):
            if text != core:
                legacy = text
                break
    description_parts = [part for part in detailed_description(event).split("\n\n") if part.strip()]
    if not core and description_parts:
        core = description_parts[0]
    if not legacy:
        legacy = description_parts[-1] if description_parts else "Archive 暂无补充说明。"
    return core, legacy


def display_year(event) -> str:
    return str(event.get("year") or event.get("date") or "时间未标注")


def display_title(event) -> str:
    presentation = event.get("defaultPresentation", {})
    return localized(presentation.get("displayTitle")) or localized(event.get("title")) or event.get("id", "")


def add_cover(
    prs,
    storyline,
    config,
    total_pages,
    event_count,
    part_index=1,
    part_count=1,
    event_position_start=None,
    event_position_end=None,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    title = deck_title(storyline)
    subtitle = localized(storyline.get("subtitle")) or localized(storyline.get("description"))
    solid_rect(slide, 0, 0, 0.22, 7.5, VERMILION)
    solid_rect(slide, 8.82, 0, 4.51, 7.5, TEAL)
    add_text(slide, config["coverTag"], 0.82, 0.76, 7.4, 0.28, 11, TEAL, True, FONT_LATIN, fit=False)
    add_text(slide, title, 0.78, 1.42, 7.35, 1.36, 31, INK, True, line_spacing=1.08)
    if part_count > 1:
        volume_label = f"第 {part_index}/{part_count} 册"
        if event_position_start is not None and event_position_end is not None:
            volume_label += f" · 事件 {event_position_start:03d}–{event_position_end:03d}"
    else:
        volume_label = "全量事件文档"
    add_text(slide, volume_label, 0.82, 3.04, 5.8, 0.48, 21, VERMILION, True)
    add_rule(slide, 0.82, 3.78, 6.86, HAIRLINE)
    add_text(slide, subtitle, 0.82, 4.08, 6.58, 1.04, 15.5, INK_SOFT, line_spacing=1.28)
    add_text(slide, "Archive JSON 驱动 · 每个事件两页 · 可编辑 PPTX", 0.82, 5.46, 6.55, 0.32, 10.5, INK_SOFT, True)
    add_text(slide, f"{event_count:03d}", 9.26, 1.08, 3.02, 1.08, 51, WHITE, True, FONT_LATIN)
    add_text(slide, "EVENTS", 9.34, 2.31, 2.10, 0.30, 13, WHITE, True, FONT_LATIN, fit=False)
    add_text(slide, "历史概览\n技术 / 主题解释\n人物与资料图\n事件页码索引", 9.34, 3.18, 2.84, 2.25, 17, WHITE, True, line_spacing=1.42)
    add_footer(slide, title, 1, total_pages)


def add_contents_page(prs, storyline, events, page_number, total_pages, toc_page_index, toc_page_count, first_event_page):
    title = deck_title(storyline)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_header(
        slide,
        title,
        f"事件目录 · {toc_page_index:02d}/{toc_page_count:02d}",
        title,
        f"共 {len(events)} 个事件",
        page_number,
        total_pages,
    )
    page_events = events[(toc_page_index - 1) * TOC_EVENTS_PER_PAGE : toc_page_index * TOC_EVENTS_PER_PAGE]
    y = 1.38
    for local_index, event in enumerate(page_events):
        local_position = (toc_page_index - 1) * TOC_EVENTS_PER_PAGE + local_index + 1
        storyline_position = event["_storylinePosition"]
        color = TEAL if (storyline_position - 1) % 2 == 0 else VERMILION
        event_page = first_event_page + (local_position - 1) * 2
        add_text(slide, f"{storyline_position:03d}", 0.73, y, 0.75, 0.32, 14, color, True, FONT_LATIN, fit=False)
        add_text(slide, display_title(event), 1.55, y - 0.01, 6.36, 0.34, 14.5, INK, True)
        add_text(slide, display_year(event), 8.15, y + 0.02, 1.34, 0.25, 10, INK_SOFT, True, FONT_LATIN)
        add_text(slide, f"第 {event_page}–{event_page + 1} 页", 10.08, y + 0.02, 2.08, 0.25, 9.5, INK_SOFT, align=PP_ALIGN.RIGHT)
        add_rule(slide, 1.54, y + 0.39, 10.62, HAIRLINE, 0.5)
        y += 0.56


def add_overview_slide(prs, storyline, event, position, page, total_pages, figure_map, cache_dir):
    storyline_title = deck_title(storyline)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    year = display_year(event)
    title = display_title(event)
    add_header(
        slide,
        storyline_title,
        f"事件 {position:03d} · 历史概览",
        title,
        year,
        page,
        total_pages,
        event.get("id"),
    )
    overview_asset, used_fallback = choose_asset(event, OVERVIEW_ROLES)
    solid_rect(slide, 0.66, 1.36, 3.72, 5.55, WHITE, HAIRLINE)
    contain = not overview_asset or overview_asset.get("role") not in FILL_ROLES
    add_picture(slide, overview_asset, cache_dir, 0.84, 1.54, 3.36, 4.24, contain=contain)
    add_text(slide, year, 0.92, 6.02, 2.35, 0.50, 24, VERMILION, True, FONT_LATIN)
    add_text(slide, "AI HISTORY", 3.02, 6.17, 1.00, 0.18, 8, INK_SOFT, True, FONT_LATIN, PP_ALIGN.RIGHT, fit=False)
    add_label(slide, "事件详述", 4.80, 1.43, 2.0)
    add_paragraph_text(slide, detailed_description(event), 4.78, 1.79, 7.72, 3.20, 16.5, INK, line_spacing=1.23, paragraph_spacing=9)
    add_rule(slide, 4.80, 5.15, 7.67, TEAL, 1.4)
    add_label(slide, "关键人物", 4.80, 5.39, 2.0)
    figures = figure_lines(event, figure_map)
    if not figures:
        figures = [(localized(event.get("location", {}).get("place"), default="相关机构"), "事件相关机构或地点")]
    row_y = 5.71
    for index, (name, role) in enumerate(figures[:4]):
        column = index % 2
        row = index // 2
        x = 4.80 + column * 3.87
        y = row_y + row * 0.62
        solid_rect(slide, x, y + 0.04, 0.07, 0.43, TEAL if index == 0 else OCHRE)
        add_text(slide, name, x + 0.17, y, 3.43, 0.25, 13, INK, True)
        add_text(slide, role, x + 0.17, y + 0.27, 3.43, 0.22, 9.5, INK_SOFT)
    topics = " / ".join(event.get("topics", [])) or storyline.get("id", "AI 历史")
    add_text(slide, topics, 4.80, 6.78, 7.55, 0.16, 8.5, TEAL, True, FONT_LATIN, fit=False)
    return overview_asset, used_fallback


def add_explanation_slide(prs, storyline, event, position, page, total_pages, cache_dir, overview_asset):
    storyline_title = deck_title(storyline)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    title = display_title(event)
    add_header(
        slide,
        storyline_title,
        f"事件 {position:03d} · 主题解释",
        title,
        "核心思想与长期影响",
        page,
        total_pages,
        event.get("id"),
    )
    excluded = [overview_asset.get("id")] if overview_asset else []
    explainer_asset, used_fallback = choose_asset(event, EXPLAINER_ROLES, excluded)
    solid_rect(slide, 0.66, 1.36, 6.35, 5.55, WHITE, HAIRLINE)
    add_picture(slide, explainer_asset, cache_dir, 0.93, 1.62, 5.81, 4.72, contain=True)
    add_text(slide, "ARCHIVE VISUAL", 0.94, 6.54, 2.25, 0.16, 7.5, INK_SOFT, True, FONT_LATIN, fit=False)
    core, legacy = content_sections(event)
    solid_rect(slide, 7.36, 1.36, 5.31, 2.60, TEAL_PALE, TEAL_PALE)
    add_label(slide, "核心思想", 7.70, 1.68, 1.5, TEAL)
    add_text(slide, core, 7.68, 2.06, 4.55, 1.48, 16.5, INK, line_spacing=1.25)
    solid_rect(slide, 7.36, 4.24, 5.31, 2.67, WHITE, HAIRLINE)
    add_label(slide, "长期影响", 7.70, 4.57, 1.5, VERMILION)
    add_text(slide, legacy, 7.68, 4.96, 4.55, 1.52, 15.5, INK_SOFT, line_spacing=1.25)
    return explainer_asset, used_fallback


def validate_layout(prs):
    issues = []
    for slide_number, slide in enumerate(prs.slides, 1):
        for shape_number, shape in enumerate(slide.shapes, 1):
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > prs.slide_width or shape.top + shape.height > prs.slide_height:
                issues.append(f"slide {slide_number} shape {shape_number} exceeds slide bounds")
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                if shape.width < Inches(0.18) or shape.height < Inches(0.12):
                    issues.append(f"slide {slide_number} text shape {shape_number} is too small")
    if issues:
        raise RuntimeError("\n".join(issues))


def generate_storyline(
    config,
    output_dir: Path,
    figure_map,
    event_references=None,
    part_index=1,
    part_count=1,
    filename=None,
    index_filename=None,
) -> dict:
    storyline = load_storyline(config["id"])
    all_references = [dict(reference) for reference in storyline["events"]]
    for position, reference in enumerate(all_references, 1):
        reference["_storylinePosition"] = position
    references = event_references or all_references
    events = []
    for reference in references:
        event = load_event(reference["eventId"])
        event["_storylinePosition"] = reference["_storylinePosition"]
        events.append(event)
    toc_pages = math.ceil(len(events) / TOC_EVENTS_PER_PAGE)
    total_pages = 1 + toc_pages + len(events) * 2
    first_event_page = 2 + toc_pages
    cache_dir = ROOT / ".tmp" / "storyline-ppt-assets"
    cache_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    add_cover(
        prs,
        storyline,
        config,
        total_pages,
        len(events),
        part_index,
        part_count,
        events[0]["_storylinePosition"],
        events[-1]["_storylinePosition"],
    )
    for toc_index in range(1, toc_pages + 1):
        add_contents_page(
            prs,
            storyline,
            events,
            1 + toc_index,
            total_pages,
            toc_index,
            toc_pages,
            first_event_page,
        )
    index = {
        "storylineId": storyline["id"],
        "title": localized(storyline.get("title")),
        "partIndex": part_index,
        "partCount": part_count,
        "sourceEventCount": len(all_references),
        "eventCount": len(events),
        "slideCount": total_pages,
        "tocSlideCount": toc_pages,
        "events": [],
    }
    overview_fallbacks = 0
    explainer_fallbacks = 0
    for local_position, event in enumerate(events, 1):
        storyline_position = event["_storylinePosition"]
        overview_page = first_event_page + (local_position - 1) * 2
        overview_asset, overview_fallback = add_overview_slide(
            prs,
            storyline,
            event,
            storyline_position,
            overview_page,
            total_pages,
            figure_map,
            cache_dir,
        )
        explainer_asset, explainer_fallback = add_explanation_slide(
            prs,
            storyline,
            event,
            storyline_position,
            overview_page + 1,
            total_pages,
            cache_dir,
            overview_asset,
        )
        overview_fallbacks += int(overview_fallback)
        explainer_fallbacks += int(explainer_fallback)
        index["events"].append(
            {
                "eventId": event["id"],
                "title": display_title(event),
                "year": display_year(event),
                "storylinePosition": storyline_position,
                "slideStart": overview_page,
                "slideEnd": overview_page + 1,
                "overviewAssetId": overview_asset.get("id") if overview_asset else None,
                "explanationAssetId": explainer_asset.get("id") if explainer_asset else None,
            }
        )
    validate_layout(prs)
    output_path = output_dir / (filename or config["filename"])
    index_path = output_dir / (index_filename or config["indexFilename"])
    prs.save(output_path)
    index_path.write_text(json.dumps(index, ensure_ascii=False, indent=2), encoding="utf-8")
    return {
        "storylineId": storyline["id"],
        "title": localized(storyline.get("title")),
        "partIndex": part_index,
        "partCount": part_count,
        "sourceEventCount": len(all_references),
        "eventPositionStart": events[0]["_storylinePosition"],
        "eventPositionEnd": events[-1]["_storylinePosition"],
        "eventCount": len(events),
        "slideCount": total_pages,
        "tocSlideCount": toc_pages,
        "overviewFallbackCount": overview_fallbacks,
        "explanationFallbackCount": explainer_fallbacks,
        "pptx": str(output_path.relative_to(ROOT)),
        "index": str(index_path.relative_to(ROOT)),
        "sizeBytes": output_path.stat().st_size,
    }


def split_balanced(items, part_count):
    base_size, remainder = divmod(len(items), part_count)
    chunks = []
    start = 0
    for part_index in range(part_count):
        size = base_size + (1 if part_index < remainder else 0)
        chunks.append(items[start : start + size])
        start += size
    return chunks


def part_output_names(config, part_index):
    prefix, separator, rest = config["filename"].partition("-")
    if not separator:
        raise ValueError(f"Cannot derive part filename from {config['filename']}")
    letter = chr(ord("a") + part_index - 1)
    stem = Path(rest).stem
    filename = f"{prefix}{letter}-{stem}-part-{part_index}.pptx"
    index_filename = f"{prefix}{letter}-{stem}-part-{part_index}-index.json"
    return filename, index_filename


def parts_manifest_name(config):
    return f"{Path(config['filename']).stem}-parts-manifest.json"


def parse_args():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--storyline",
        action="append",
        choices=[config["id"] for config in STORYLINES],
        help="Generate only the selected storyline; repeat for multiple storylines.",
    )
    parser.add_argument(
        "--parts",
        type=int,
        default=1,
        help="Split each selected storyline into this many balanced, contiguous decks.",
    )
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    args = parser.parse_args()
    if args.parts < 1 or args.parts > 26:
        parser.error("--parts must be between 1 and 26")
    return args


def main():
    args = parse_args()
    output_dir = args.output if args.output.is_absolute() else ROOT / args.output
    output_dir.mkdir(parents=True, exist_ok=True)
    selected = set(args.storyline or [config["id"] for config in STORYLINES])
    figure_map = load_figures()
    manifest_path = output_dir / "manifest.json"
    existing_manifest = []
    if manifest_path.exists():
        existing_manifest = load_json(manifest_path)
    reports_by_id = {report["storylineId"]: report for report in existing_manifest}
    for config in STORYLINES:
        if config["id"] not in selected:
            continue
        if args.parts > 1:
            storyline = load_storyline(config["id"])
            references = [dict(reference) for reference in storyline["events"]]
            if args.parts > len(references):
                raise ValueError(
                    f"Cannot split {len(references)} events into {args.parts} non-empty parts"
                )
            for position, reference in enumerate(references, 1):
                reference["_storylinePosition"] = position
            part_reports = []
            for part_index, chunk in enumerate(split_balanced(references, args.parts), 1):
                filename, index_filename = part_output_names(config, part_index)
                print(
                    f"Generating {config['id']} part {part_index}/{args.parts}...",
                    flush=True,
                )
                report = generate_storyline(
                    config,
                    output_dir,
                    figure_map,
                    event_references=chunk,
                    part_index=part_index,
                    part_count=args.parts,
                    filename=filename,
                    index_filename=index_filename,
                )
                part_reports.append(report)
                print(
                    f"Generated {report['pptx']}: {report['eventCount']} events, "
                    f"{report['slideCount']} slides, {report['sizeBytes']} bytes",
                    flush=True,
                )
            parts_manifest_path = output_dir / parts_manifest_name(config)
            parts_manifest_path.write_text(
                json.dumps(part_reports, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
            print(parts_manifest_path, flush=True)
            continue
        print(f"Generating {config['id']}...", flush=True)
        report = generate_storyline(config, output_dir, figure_map)
        reports_by_id[report["storylineId"]] = report
        print(
            f"Generated {report['pptx']}: {report['eventCount']} events, "
            f"{report['slideCount']} slides, {report['sizeBytes']} bytes",
            flush=True,
        )
    if args.parts > 1:
        return
    manifest = [reports_by_id[config["id"]] for config in STORYLINES if config["id"] in reports_by_id]
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    print(manifest_path, flush=True)


if __name__ == "__main__":
    main()
